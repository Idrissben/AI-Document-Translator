#type: ignore
import io
import pytest
from django.core.files.uploadedfile import SimpleUploadedFile
from django.urls import reverse
from django.test import Client
import docx
from pptx import Presentation
from openpyxl import Workbook


@pytest.mark.django_db
def test_upload_and_translate_docx():
    client = Client()

    # Create an in-memory DOCX file
    doc_file = io.BytesIO()
    doc = docx.Document()
    doc.add_paragraph("sample content")
    doc.save(doc_file)
    doc_file.seek(0)
    uploaded_file = SimpleUploadedFile(
        "test.docx",
        doc_file.read(),
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )

    response = client.post(
        reverse("upload_and_translate"),
        {
            "document": uploaded_file,
            "source_language": "fr",
            "target_language": "en",
            "evaluation_method": "no_evaluation",
        },
    )

    assert response.status_code == 200  # Expect successful translation


@pytest.mark.django_db
def test_upload_and_translate_pptx():
    client = Client()

    # Create an in-memory PPTX file
    pptx_file = io.BytesIO()
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "This is a powerpoint"
    prs.save(pptx_file)
    pptx_file.seek(0)
    uploaded_file = SimpleUploadedFile(
        "test.pptx",
        pptx_file.read(),
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    response = client.post(
        reverse("upload_and_translate"),
        {
            "document": uploaded_file,
            "source_language": "fr",
            "target_language": "en",
            "evaluation_method": "no_evaluation",
        },
    )

    assert response.status_code == 200  # Expect successful translation


@pytest.mark.django_db
def test_read_write_excel(tmp_path):
    # Create and save an Excel file in-memory
    excel_file = io.BytesIO()
    wb = Workbook()
    ws = wb.active
    ws.append(["Col1", "Col2"])
    ws.append(["Val1", "Val2"])
    wb.save(excel_file)
    excel_file.seek(0)
    uploaded_file = SimpleUploadedFile(
        "test.xlsx",
        excel_file.read(),
        content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )

    client = Client()
    response = client.post(
        reverse("upload_and_translate"),
        {
            "document": uploaded_file,
            "source_language": "fr",
            "target_language": "en",
            "evaluation_method": "no_evaluation",
        },
    )

    assert response.status_code == 200  # Expect successful translation
