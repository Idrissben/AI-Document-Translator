# type: ignore
from .utils import read_docx, write_docx, read_excel, write_excel, read_pptx, write_pptx
from .translator import initialize_model, translate_text
from langfuse import Langfuse
from pptx import Presentation
import os
import io

model = initialize_model()

trace = Langfuse()

# Sample data
sample_text = "Bonjour tout le monde"
reference_translation = "Hello everyone"
reference_translation_glossary = "Hello Casablanca"
glossary = {"tout le monde": "Casablanca"}


def test_translate_text_no_glossary():
    result = translate_text(sample_text, model, trace=trace)
    assert (
        result == reference_translation
    ), "Translation result should match the expected translation"


def test_translate_text_with_glossary():
    result = translate_text(sample_text, model, glossary=glossary, trace=trace)
    assert (
        result == reference_translation_glossary
    ), "Translation result should match the expected translation"


# Test read and write docx functions
def test_read_write_docx(tmp_path):
    docx_file = tmp_path / "test.docx"
    write_docx(sample_text, str(docx_file))
    assert os.path.exists(docx_file), "DOCX file should be created"

    read_text = read_docx(str(docx_file))
    assert read_text == sample_text, "Read text should match the written text"


def test_read_write_pptx(tmp_path):
    sample_text = "This is a sample text"
    slides_text = [[(0, sample_text)]]
    pptx_output_file = tmp_path / "test_output.pptx"

    template_pptx = io.BytesIO()
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])  # Add a single slide
    prs.save(template_pptx)
    template_pptx.seek(0)  # Reset the file pointer to the start

    with open(pptx_output_file, "wb") as output_file:
        output_file.write(template_pptx.read())  # Copy the template contents to output

    # Call `write_pptx` function with the in-memory template path
    write_pptx(slides_text, str(pptx_output_file), str(pptx_output_file))
    read_slides = read_pptx(str(pptx_output_file))

    assert isinstance(read_slides, list), "Read PPTX content should be a list of slides"
    assert len(read_slides) == len(
        slides_text
    ), "Number of slides should match the input"
    assert (
        read_slides[0][0][1] == sample_text
    ), "Read text should match the written text on the first slide"


# Test read and write Excel functions
def test_read_write_excel(tmp_path):
    excel_file = tmp_path / "test.xlsx"
    write_excel("Col1\tCol2\nVal1\tVal2", str(excel_file))
    assert os.path.exists(excel_file), "Excel file should be created"

    read_text = read_excel(str(excel_file))
    assert "Col1\tCol2" in read_text, "Read content should include column headers"
    assert "Val1\tVal2" in read_text, "Read content should include cell values"
