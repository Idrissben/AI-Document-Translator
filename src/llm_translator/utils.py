# type: ignore
import docx
from pptx import Presentation
import openpyxl
import nltk
from nltk.translate.meteor_score import meteor_score

try:
    nltk.data.find("corpora/wordnet")
except LookupError:
    nltk.download("wordnet")


# Functions to read and write .docx files
def read_docx(file_path: str):
    doc = docx.Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip() != "":
            full_text.append(para.text)
    return "\n".join(full_text)


def write_docx(text: str, file_path: str):
    doc = docx.Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    doc.save(file_path)


def read_pptx(file_path: str):
    prs = Presentation(file_path)
    slides_content = []
    for slide in prs.slides:
        slide_texts = []
        for idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip() != "":
                slide_texts.append((idx, shape.text))  # Store index and text
        slides_content.append(
            slide_texts
        )  # List of (shape_index, text) tuples per slide
    return slides_content  # Returns a list of lists


def write_pptx(slides_texts: list, file_path: str, template_path: str):
    prs = Presentation(template_path)
    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx < len(slides_texts):
            texts = slides_texts[
                slide_idx
            ]  # List of (shape_index, translated_text) tuples
            for shape_idx, translated_text in texts:
                shape = slide.shapes[shape_idx]
                if shape.has_text_frame:
                    shape.text_frame.clear()
                    shape.text_frame.text = translated_text
    prs.save(file_path)


# Functions to read and write Excel (.xlsx) files
def read_excel(file_path: str):
    workbook = openpyxl.load_workbook(file_path)
    full_text = []
    for sheet in workbook.sheetnames:
        ws = workbook[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = [str(cell) if cell is not None else "" for cell in row]
            full_text.append("\t".join(row_text))
    return "\n".join(full_text)


def write_excel(text: str, file_path: str):
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    for line in text.split("\n"):
        row_data = line.split("\t")
        sheet.append(row_data)
    workbook.save(file_path)


def calculate_metric(reference_translation: str, predicted_translation: str):
    return meteor_score([reference_translation.split()], predicted_translation.split())
