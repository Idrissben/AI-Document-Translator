#type: ignore
import docx
from pptx import Presentation
import openpyxl
import nltk
from nltk.translate import meteor_score

try:
    nltk.data.find('corpora/wordnet')
except LookupError:
    nltk.download('wordnet')

# Functions to read and write .docx files
def read_docx(file_path: str):
    doc = docx.Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip() != "":
            full_text.append(para.text)
    return "\n".join(full_text)


def write_docx(text: str, file_path: str):
    doc = docx.Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    doc.save(file_path)


# Function to read .pptx files and preserve slide structure
def read_pptx(file_path: str):
    prs = Presentation(file_path)
    slides_content = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() != "":
                slide_text.append(shape.text)
        slides_content.append("\n".join(slide_text))  # Combine text per slide
    return slides_content  # Return a list with each slide's text separately


# Function to write .pptx files with the same slide structure
def write_pptx(slides_text: str, file_path: str, template_path: str):
    prs = Presentation(template_path)

    for idx, slide in enumerate(prs.slides):
        # Ensure that we have enough slides in input text
        if idx < len(slides_text):
            text = slides_text[idx]
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape.text = text  # Replace with translated text for each slide

    prs.save(file_path)


# Functions to read and write Excel (.xlsx) files
def read_excel(file_path: str):
    workbook = openpyxl.load_workbook(file_path)
    full_text = []
    for sheet in workbook.sheetnames:
        ws = workbook[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = [str(cell) if cell is not None else "" for cell in row]
            full_text.append("\t".join(row_text))
    return "\n".join(full_text)


def write_excel(text: str, file_path: str):
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    for line in text.split("\n"):
        row_data = line.split("\t")
        sheet.append(row_data)
    workbook.save(file_path)

def calculate_metric(reference_translation: str, predicted_translation: str):
    return meteor_score([reference_translation.split()], predicted_translation.split())